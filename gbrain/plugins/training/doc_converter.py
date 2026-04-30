"""
文档转换模块 - 支持 PDF/DOCX/PPT/TXT 转换为 Markdown（增强版）
"""

import os
import re
from pathlib import Path
from typing import Optional

import aiofiles


class DocumentConverter:
    """文档格式转换器"""

    SUPPORTED_FORMATS = ['.pdf', '.docx', '.pptx', '.txt', '.md']

    def __init__(self, upload_dir: str = None):
        from gbrain.config import BASE_PATH
        if upload_dir is None:
            upload_dir = BASE_PATH / "data" / "uploads"
        self.upload_dir = Path(upload_dir)
        self.upload_dir.mkdir(parents=True, exist_ok=True)

    async def save_file(self, filename: str, content: bytes) -> Path:
        """保存上传文件"""
        file_path = self.upload_dir / filename
        async with aiofiles.open(file_path, 'wb') as f:
            await f.write(content)
        return file_path

    def convert_to_markdown(self, file_path: Path) -> str:
        """根据文件类型转换为 Markdown"""
        ext = file_path.suffix.lower()

        if ext == '.pdf':
            return self._convert_pdf_to_markdown(file_path)
        elif ext == '.docx':
            return self._convert_docx_to_markdown(file_path)
        elif ext == '.pptx':
            return self._convert_pptx_to_markdown(file_path)
        elif ext == '.txt':
            return self._convert_txt_to_markdown(file_path)
        elif ext == '.md':
            return self._convert_md_to_markdown(file_path)
        else:
            raise ValueError(f"不支持的文件格式: {ext}")

    def _clean_text(self, text: str) -> str:
        """清理文本：去除多余空白、规范化换行"""
        if not text:
            return ""
        # 去除多余空格（保留单词间单个空格）
        text = re.sub(r'[ \t]+', ' ', text)
        # 去除行首行尾空白
        text = text.strip()
        return text

    def _normalize_whitespace(self, text: str) -> str:
        """规范化连续空白字符为单个换行或空格"""
        # 将3个以上连续换行转为2个
        text = re.sub(r'\n{3,}', '\n\n', text)
        # 将2个以上连续空格+换行组合处理
        text = re.sub(r' +\n', '\n', text)
        return text

    def _convert_pdf_to_markdown(self, file_path: Path) -> str:
        """PDF 转 Markdown（增强版：使用 pdfplumber 更好保留结构）"""
        try:
            import pdfplumber
        except ImportError:
            return self._convert_pdf_fallback(file_path)

        markdown_parts = []

        with pdfplumber.open(str(file_path)) as pdf:
            for i, page in enumerate(pdf.pages):
                # 提取文本，保留段落结构
                text = page.extract_text()

                if not text or not text.strip():
                    continue

                # 尝试提取表格
                tables = page.extract_tables()
                table_text = ""
                if tables:
                    for table in tables:
                        if table:
                            table_lines = []
                            for row in table:
                                # 过滤 None 值
                                cells = [str(cell).strip() if cell else "" for cell in row]
                                # 移除空行
                                cells = [c for c in cells if c]
                                if cells:
                                    table_lines.append("| " + " | ".join(cells) + " |")
                            if table_lines:
                                table_text += "\n" + "\n".join(table_lines) + "\n"

                # 清理文本
                text = self._clean_text(text)

                # 将连续短行合并为段落（保留标题结构）
                lines = text.split('\n')
                paragraph_lines = []
                current_paragraph = []

                for line in lines:
                    line = line.strip()
                    if not line:
                        # 遇到空行，合并当前段落
                        if current_paragraph:
                            para_text = ' '.join(current_paragraph)
                            if para_text:
                                paragraph_lines.append(para_text)
                            current_paragraph = []
                        continue

                    # 检查是否是标题（短行+常见标题模式）
                    is_heading = (
                        len(line) < 100 and
                        (
                            re.match(r'^(第[一二三四五六七八九十\d]+[章节页]|[一二三四五六七八九十百千\d]+[.、)）])', line) or
                            re.match(r'^(【|『|\[)', line) or
                            re.match(r'^\d+[.、]\s*[\u4e00-\u9fa5]', line) or
                            re.match(r'^[\u4e00-\u9fa5]{2,6}(:|：|$)', line) or
                            line.endswith('：') or
                            line.endswith(':')
                        )
                    )

                    if is_heading and len(line) < 50:
                        # 先输出当前段落
                        if current_paragraph:
                            para_text = ' '.join(current_paragraph)
                            if para_text:
                                paragraph_lines.append(para_text)
                            current_paragraph = []
                        paragraph_lines.append(line)
                    else:
                        current_paragraph.append(line)

                # 处理最后一段
                if current_paragraph:
                    para_text = ' '.join(current_paragraph)
                    if para_text:
                        paragraph_lines.append(para_text)

                page_content = '\n\n'.join(paragraph_lines)

                # 如果有表格，追加到内容后面
                if table_text:
                    page_content += "\n\n" + table_text

                if page_content.strip():
                    markdown_parts.append(f"## 第 {i + 1} 页\n\n{page_content.strip()}")

        result = '\n\n'.join(markdown_parts)
        return self._normalize_whitespace(result)

    def _convert_pdf_fallback(self, file_path: Path) -> str:
        """PDF 回退方案（使用 pypdf）"""
        try:
            from pypdf import PdfReader
        except ImportError:
            # 极简回退：直接读取二进制
            with open(file_path, 'rb') as f:
                content = f.read()
                try:
                    return content.decode('utf-8')
                except:
                    return content.decode('gbk', errors='ignore')

        reader = PdfReader(str(file_path))
        markdown_parts = []

        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text and text.strip():
                text = self._clean_text(text)
                markdown_parts.append(f"## 第 {i + 1} 页\n\n{text.strip()}")

        result = '\n\n'.join(markdown_parts)
        return self._normalize_whitespace(result)

    def _convert_docx_to_markdown(self, file_path: Path) -> str:
        """DOCX 转 Markdown（增强版：更好保留 inline 样式和结构）"""
        from docx import Document
        from docx.shared import Pt
        from docx.enum.text import WD_PARAGRAPH_ALIGNMENT

        doc = Document(str(file_path))
        markdown_parts = []

        def get_paragraph_markdown(para) -> str:
            """将段落转换为 Markdown，支持 inline 样式"""
            # 获取段落样式
            style_name = para.style.name.lower() if para.style else ''

            # 检查标题级别
            if 'heading 1' in style_name or style_name == 'title':
                return f"# {get_run_text(para)}"
            elif 'heading 2' in style_name:
                return f"## {get_run_text(para)}"
            elif 'heading 3' in style_name:
                return f"### {get_run_text(para)}"
            elif 'heading 4' in style_name:
                return f"#### {get_run_text(para)}"

            # 检查列表
            if 'list' in style_name or 'bullet' in style_name:
                return f"- {get_run_text(para)}"

            # 普通段落
            text = get_run_text_with_formatting(para)
            return text if text.strip() else None

        def get_run_text(para) -> str:
            """获取段落纯文本"""
            return ''.join(run.text for run in para.runs)

        def get_run_text_with_formatting(para) -> str:
            """获取段落文本，保留粗体/斜体标记"""
            parts = []
            for run in para.runs:
                text = run.text
                if not text:
                    continue

                # 检查粗体
                if run.bold:
                    # 检查是否已经是 bold 标记内容
                    if text.startswith('**') and text.endswith('**'):
                        parts.append(text)
                    else:
                        parts.append(f"**{text}**")
                elif run.italic:
                    if text.startswith('*') and text.endswith('*'):
                        parts.append(text)
                    else:
                        parts.append(f"*{text}*")
                else:
                    parts.append(text)

            result = ''.join(parts)
            return result

        # 处理段落
        for para in doc.paragraphs:
            result = get_paragraph_markdown(para)
            if result:
                markdown_parts.append(result)

        # 处理表格
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                cells = []
                for cell in row.cells:
                    # 清理单元格文本
                    cell_text = cell.text.strip()
                    cell_text = re.sub(r'\s+', ' ', cell_text)
                    cells.append(cell_text)
                if cells:
                    table_text.append("| " + " | ".join(cells) + " |")

            if table_text:
                # 检查表头分隔行
                if table_text:
                    markdown_parts.append("\n".join(table_text))
                    # 添加分隔行（根据列数）
                    if len(table_text) > 0:
                        col_count = table_text[0].count('|') - 1
                        markdown_parts.append("| " + " | ".join(["---"] * col_count) + " |")

        result = '\n\n'.join(markdown_parts)
        return self._normalize_whitespace(result)

    def _convert_pptx_to_markdown(self, file_path: Path) -> str:
        """PPTX 转 Markdown（增强版：更好提取结构和表格）"""
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation(str(file_path))
        markdown_parts = []

        for i, slide in enumerate(prs.slides):
            slide_title = ""
            slide_content = []
            slide_tables = []

            # 获取所有形状
            for shape in slide.shapes:
                # 处理表格
                if shape.has_table:
                    table = shape.table
                    table_text = []
                    for row_idx, row in enumerate(table.rows):
                        cells = []
                        for cell in row.cells:
                            cell_text = cell.text_frame.text.strip()
                            cell_text = re.sub(r'\s+', ' ', cell_text)
                            cells.append(cell_text)
                        table_text.append("| " + " | ".join(cells) + " |")

                        # 添加分隔行
                        if row_idx == 0:
                            col_count = len(cells)
                            table_text.append("| " + " | ".join(["---"] * col_count) + " |")

                    if table_text:
                        slide_tables.append('\n'.join(table_text))

                # 处理文本
                if hasattr(shape, "text") and shape.has_text_frame:
                    text = shape.text.strip()
                    if not text:
                        continue

                    # 判断是否是标题
                    is_title = False
                    if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                        ph_type = str(shape.placeholder_format.type) if shape.placeholder_format else ""
                        if "TITLE" in ph_type or "CENTER_TITLE" in ph_type:
                            is_title = True
                            slide_title = f"## {text}"
                            continue

                    # 检查形状位置（顶部区域可能是标题）
                    if hasattr(shape, "top") and shape.top:
                        top_inches = shape.top / 914400  # EMU to inches
                        if top_inches < 1.5:  # 顶部 1.5 英寸内
                            # 检查字体大小
                            if shape.text_frame.paragraphs:
                                first_para = shape.text_frame.paragraphs[0]
                                if first_para.runs:
                                    for run in first_para.runs:
                                        if hasattr(run, "font_size") and run.font_size:
                                            if run.font_size >= Pt(24):
                                                is_title = True
                                                slide_title = f"## {text}"
                                                continue

                    if not is_title:
                        slide_content.append(text)

            # 组合幻灯片内容
            if slide_title or slide_content or slide_tables:
                markdown_parts.append(f"## 幻灯片 {i + 1}")

                if slide_title:
                    markdown_parts.append(slide_title)

                # 添加要点列表
                if slide_content:
                    content_items = []
                    for item in slide_content:
                        # 跳过太短的（可能是标题重复）
                        if len(item) > 3:
                            content_items.append(f"- {item}")
                    if content_items:
                        markdown_parts.append('\n'.join(content_items))

                # 添加表格
                if slide_tables:
                    markdown_parts.append('\n\n'.join(slide_tables))

                markdown_parts.append("")  # 空行分隔

        result = '\n'.join(markdown_parts)
        return self._normalize_whitespace(result)

    def _convert_txt_to_markdown(self, file_path: Path) -> str:
        """TXT 转 Markdown（增强版：智能识别标题和结构）"""
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()

        lines = content.split('\n')
        result_lines = []
        i = 0

        while i < len(lines):
            line = lines[i].strip()

            # 跳过空行
            if not line:
                i += 1
                continue

            # 识别标题模式
            is_heading = False

            # 模式1: 以 # 开头（已是 Markdown）
            if line.startswith('#'):
                is_heading = True

            # 模式2: 数字编号开头（1. 2. 2.1）
            elif re.match(r'^\d+(\.\d+)*[.、)\s]', line):
                if len(line) < 100:  # 标题不会太长
                    is_heading = True

            # 模式3: 特殊符号开头（【 】『 』）
            elif re.match(r'^【|『|\[', line):
                is_heading = True

            # 模式4: 特定中文字符结尾（： --）
            elif re.match(r'^[\u4e00-\u9fa5]{2,6}[:：]', line):
                is_heading = True

            # 模式5: 全大写或特定关键词
            elif line.isupper() and len(line) < 50:
                is_heading = True

            # 模式6: 短行 + 特殊分隔符（===== ----）
            elif i + 1 < len(lines) and lines[i + 1].strip() and re.match(r'^[=~-]{3,}$', lines[i + 1].strip()):
                is_heading = True

            if is_heading:
                # 清理标题格式
                clean_line = re.sub(r'^[=~\s]+|[=~\s]+$', '', line)
                if clean_line.startswith('#'):
                    result_lines.append(clean_line)
                else:
                    result_lines.append(f"### {clean_line}")
                i += 2  # 跳过标题和分隔线
                continue

            # 模式7: 列表项
            if re.match(r'^[-*•]\s', line) or re.match(r'^\d+[.、]\s', line):
                result_lines.append(line)

            # 模式8: 多行合并为段落
            else:
                # 收集连续非空行
                paragraph_lines = [line]
                j = i + 1
                while j < len(lines) and lines[j].strip():
                    next_line = lines[j].strip()
                    # 如果下一行是标题模式或列表，停止
                    if next_line.startswith('#') or re.match(r'^[-*•]\s', next_line) or re.match(r'^\d+[.、]\s', next_line):
                        break
                    paragraph_lines.append(next_line)
                    j += 1

                # 合并为段落
                paragraph = ' '.join(paragraph_lines)
                if paragraph:
                    result_lines.append(paragraph)

                i = j
                continue

            i += 1

        result = '\n\n'.join(result_lines)
        # 再次规范化空白
        result = self._normalize_whitespace(result)
        return result

    def _convert_md_to_markdown(self, file_path: Path) -> str:
        """Markdown 文件直接读取"""
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        return self._normalize_whitespace(content)

    def extract_title_from_content(self, content: str) -> str:
        """从内容中提取标题（跳过幻灯片编号类标题）"""
        lines = content.split('\n')
        for line in lines:
            line = line.strip()
            if line.startswith('#'):
                # 跳过纯幻灯片编号标题（如 "## 第 1 页" 或 "## 幻灯片 1"）
                heading_text = line.lstrip('#').strip()
                # 匹配 "第 X 页" 或 "幻灯片 X" 格式
                if re.match(r'^(第\s*\d+\s*页|幻灯片\s*\d+)$', heading_text):
                    continue
                return heading_text
        # 回退逻辑：取第一行内容，跳过幻灯片编号
        first_line = content.split('\n')[0].strip()
        if re.match(r'^(第\s*\d+\s*页|幻灯片\s*\d+)$', first_line.replace('#', '').strip()):
            # 尝试从后续行找有效内容
            for line in content.split('\n')[1:10]:
                line = line.strip()
                if line and not re.match(r'^(第\s*\d+\s*页|幻灯片\s*\d+)$', line.replace('#', '').strip()):
                    # 跳过常见的元数据行
                    if line not in ['社外秘', '人事总务部'] and not line.startswith('【'):
                        return line[:50]
        return first_line


# 全局实例
_converter: Optional[DocumentConverter] = None


def get_document_converter() -> DocumentConverter:
    global _converter
    if _converter is None:
        _converter = DocumentConverter()
    return _converter
